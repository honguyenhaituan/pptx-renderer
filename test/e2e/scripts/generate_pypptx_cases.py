#!/usr/bin/env python3
"""Generate ground truth cases using python-pptx (PPTX) + PowerPoint COM (PDF export).

Produces oracle-pypptx-* cases covering:
  - Rich text: fonts, sizes, bold/italic, alignment, vertical text, bullets
  - Shape adjustment variants: same shape with different adj values
  - Chart data variants: 2D chart types with custom data/series
  - Composite: multiple components on a single slide

Usage (from test/e2e/):
  pip install python-pptx
  python scripts/generate_pypptx_cases.py            # generate PPTX + PDF (Windows)
  python scripts/generate_pypptx_cases.py --pptx-only # generate PPTX only (any platform)
"""
from __future__ import annotations

import argparse
import json
import math
import posixpath
import random
import sys
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from pptx import Presentation
from pptx.chart.data import BubbleChartData, CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

E2E_DIR = Path(__file__).resolve().parents[1]
if str(E2E_DIR) not in sys.path:
    sys.path.insert(0, str(E2E_DIR))

CASES_DIR = E2E_DIR / "oracle" / "cases-pypptx"
TESTDATA_DIR = E2E_DIR / "testdata"
REPORT_PATH = E2E_DIR / "reports" / "oracle-failures" / "pypptx-ground-truth.json"

# Slide dimensions (standard widescreen 13.333" x 7.5")
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Case definition helpers
# ---------------------------------------------------------------------------

CaseDef = dict  # {name: str, build_fn: Callable[[Presentation], None]}

PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PML_REL_PREFIX = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS = {"p": PML_NS, "pr": REL_NS}


def _emu(inches: float) -> int:
    return int(Inches(inches))


def _rels_path(part_name: str) -> str:
    directory, filename = posixpath.split(part_name)
    return posixpath.join(directory, "_rels", f"{filename}.rels")


def _resolve_part_target(part_name: str, target: str) -> str:
    return posixpath.normpath(posixpath.join(posixpath.dirname(part_name), target))


def _relationship_target(zf: ZipFile, part_name: str, rel_suffix: str) -> str | None:
    root = etree.fromstring(zf.read(_rels_path(part_name)))
    for rel in root.xpath(".//pr:Relationship", namespaces=NS):
        rel_type = rel.get("Type", "")
        target = rel.get("Target")
        if target and rel_type == f"{PML_REL_PREFIX}{rel_suffix}":
            return _resolve_part_target(part_name, target)
    return None


def _patch_placeholder_idx_inheritance_case(pptx_path: Path) -> None:
    """Make slide placeholders inherit type through idx, matching real OOXML edge cases."""
    slide_part = "ppt/slides/slide1.xml"

    with ZipFile(pptx_path, "r") as zf:
        entries = [(info, zf.read(info.filename)) for info in zf.infolist()]
        layout_part = _relationship_target(zf, slide_part, "/slideLayout")
        if layout_part is None:
            raise RuntimeError("placeholder inheritance case slide has no slideLayout relationship")

        slide_root = etree.fromstring(dict((info.filename, data) for info, data in entries)[slide_part])
        layout_root = etree.fromstring(dict((info.filename, data) for info, data in entries)[layout_part])

    slide_title_ph = slide_root.xpath(".//p:ph[@type='title']", namespaces=NS)
    if not slide_title_ph:
        raise RuntimeError("placeholder inheritance case has no slide title placeholder")
    slide_title_ph[0].attrib.pop("type", None)
    slide_title_ph[0].set("idx", "0")

    slide_body_ph = slide_root.xpath(".//p:ph[@idx='1']", namespaces=NS)
    if not slide_body_ph:
        raise RuntimeError("placeholder inheritance case has no slide body placeholder")
    slide_body_ph[0].attrib.pop("type", None)

    layout_title_ph = layout_root.xpath(".//p:ph[@type='title']", namespaces=NS)
    if not layout_title_ph:
        raise RuntimeError("placeholder inheritance case has no layout title placeholder")
    layout_title_ph[0].set("idx", "0")

    patched = {
        slide_part: etree.tostring(slide_root, encoding="UTF-8", xml_declaration=True),
        layout_part: etree.tostring(layout_root, encoding="UTF-8", xml_declaration=True),
    }
    tmp_path = pptx_path.with_name(f"{pptx_path.name}.tmp")
    if tmp_path.exists():
        tmp_path.unlink()

    with ZipFile(tmp_path, "w", ZIP_DEFLATED) as out:
        for info, data in entries:
            out.writestr(info, patched.get(info.filename, data))
    tmp_path.replace(pptx_path)


# ---------------------------------------------------------------------------
# P0: Rich text cases
# ---------------------------------------------------------------------------

def _build_text_cases() -> list[CaseDef]:
    cases: list[CaseDef] = []
    seq = 0

    def _add(slug: str, build_fn, postprocess_fn=None):
        nonlocal seq
        seq += 1
        case = {
            "name": f"oracle-pypptx-text-{seq:04d}-{slug}",
            "build_fn": build_fn,
        }
        if postprocess_fn is not None:
            case["postprocess_fn"] = postprocess_fn
        cases.append(case)

    # --- Font families ---
    font_families = [
        ("Arial", "arial"),
        ("Times New Roman", "times-new-roman"),
        ("Calibri", "calibri"),
        ("Courier New", "courier-new"),
        ("Georgia", "georgia"),
        ("Verdana", "verdana"),
        ("Impact", "impact"),
        ("Comic Sans MS", "comic-sans"),
    ]
    for font_name, slug in font_families:
        def _build(prs, _fn=font_name):
            sld = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            txbox = sld.shapes.add_textbox(_emu(1), _emu(1), _emu(8), _emu(2))
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"The quick brown fox jumps over the lazy dog — {_fn}"
            run.font.name = _fn
            run.font.size = Pt(28)
        _add(f"font-{slug}", _build)

    # --- Font sizes ---
    for pt_size in [10, 14, 18, 24, 36, 48, 72]:
        def _build(prs, _sz=pt_size):
            sld = prs.slides.add_slide(prs.slide_layouts[6])
            txbox = sld.shapes.add_textbox(_emu(1), _emu(1), _emu(10), _emu(3))
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"Font size {_sz}pt sample text"
            run.font.size = Pt(_sz)
            run.font.name = "Calibri"
        _add(f"size-{pt_size}pt", _build)

    # --- Bold / Italic / Underline combos ---
    style_combos = [
        ("bold", True, False, False),
        ("italic", False, True, False),
        ("underline", False, False, True),
        ("bold-italic", True, True, False),
        ("bold-underline", True, False, True),
        ("bold-italic-underline", True, True, True),
    ]
    for slug, bold, italic, underline in style_combos:
        def _build(prs, _b=bold, _i=italic, _u=underline, _s=slug):
            sld = prs.slides.add_slide(prs.slide_layouts[6])
            txbox = sld.shapes.add_textbox(_emu(1), _emu(1), _emu(8), _emu(2))
            tf = txbox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"Style: {_s} — The quick brown fox"
            run.font.name = "Calibri"
            run.font.size = Pt(24)
            run.font.bold = _b
            run.font.italic = _i
            run.font.underline = _u
        _add(f"style-{slug}", _build)

    # --- Alignment ---
    alignments = [
        ("left", PP_ALIGN.LEFT),
        ("center", PP_ALIGN.CENTER),
        ("right", PP_ALIGN.RIGHT),
        ("justify", PP_ALIGN.JUSTIFY),
    ]
    for slug, align in alignments:
        def _build(prs, _a=align, _s=slug):
            sld = prs.slides.add_slide(prs.slide_layouts[6])
            txbox = sld.shapes.add_textbox(_emu(1), _emu(1), _emu(8), _emu(4))
            tf = txbox.text_frame
            tf.word_wrap = True
            for i in range(3):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"Paragraph {i+1} aligned {_s}. Lorem ipsum dolor sit amet."
                p.alignment = _a
                p.font.size = Pt(18)
                p.font.name = "Calibri"
        _add(f"align-{slug}", _build)

    # --- Font colors ---
    colors = [
        ("red", RGBColor(0xFF, 0x00, 0x00)),
        ("green", RGBColor(0x00, 0xB0, 0x50)),
        ("blue", RGBColor(0x00, 0x70, 0xC0)),
        ("orange", RGBColor(0xFF, 0xC0, 0x00)),
        ("purple", RGBColor(0x7B, 0x2D, 0x8E)),
    ]
    for slug, color in colors:
        def _build(prs, _c=color, _s=slug):
            sld = prs.slides.add_slide(prs.slide_layouts[6])
            txbox = sld.shapes.add_textbox(_emu(1), _emu(1), _emu(8), _emu(2))
            tf = txbox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"Color: {_s} text sample"
            run.font.name = "Calibri"
            run.font.size = Pt(28)
            run.font.color.rgb = _c
        _add(f"color-{slug}", _build)

    # --- Multi-paragraph with mixed formatting ---
    def _build_mixed(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = sld.shapes.add_textbox(_emu(1), _emu(0.5), _emu(10), _emu(6))
        tf = txbox.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = "Title in Bold 36pt"
        r1.font.name = "Arial"
        r1.font.size = Pt(36)
        r1.font.bold = True

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = "Subtitle in italic 24pt — "
        r2.font.name = "Georgia"
        r2.font.size = Pt(24)
        r2.font.italic = True
        r3 = p2.add_run()
        r3.text = "with colored segment"
        r3.font.name = "Georgia"
        r3.font.size = Pt(24)
        r3.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)

        p3 = tf.add_paragraph()
        r4 = p3.add_run()
        r4.text = "Body text in Calibri 18pt. Lorem ipsum dolor sit amet, consectetur adipiscing elit."
        r4.font.name = "Calibri"
        r4.font.size = Pt(18)
    _add("mixed-formatting", _build_mixed)

    # --- Bullet list ---
    def _build_bullets(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = sld.shapes.add_textbox(_emu(1), _emu(1), _emu(8), _emu(5))
        tf = txbox.text_frame
        tf.word_wrap = True
        items = [
            (0, "First level item A"),
            (1, "Second level item A.1"),
            (1, "Second level item A.2"),
            (0, "First level item B"),
            (1, "Second level item B.1"),
            (2, "Third level item B.1.a"),
            (0, "First level item C"),
        ]
        for i, (level, text) in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.level = level
            p.font.size = Pt(18)
            p.font.name = "Calibri"
    _add("bullet-list", _build_bullets)

    # --- Vertical text (East Asian) ---
    def _build_vertical(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = sld.shapes.add_textbox(_emu(4), _emu(0.5), _emu(3), _emu(6))
        tf = txbox.text_frame
        tf.word_wrap = True
        txBody = tf._txBody
        bodyPr = txBody.find(qn("a:bodyPr"))
        bodyPr.set("vert", "eaVert")
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "垂直文本テスト Vertical Text 수직 텍스트"
        run.font.size = Pt(24)
        run.font.name = "Microsoft YaHei"
    _add("vertical-east-asian", _build_vertical)

    # --- Vertical text (wordArtVert) ---
    def _build_vertical_word(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = sld.shapes.add_textbox(_emu(4), _emu(0.5), _emu(2), _emu(6))
        tf = txbox.text_frame
        tf.word_wrap = True
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        bodyPr.set("vert", "wordArtVert")
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "VERTICAL STACKED TEXT"
        run.font.size = Pt(24)
        run.font.name = "Arial"
    _add("vertical-stacked", _build_vertical_word)

    # --- Text anchor (top / middle / bottom) ---
    _ANCHOR_XML = {"top": "t", "middle": "ctr", "bottom": "b"}
    for anchor_slug in ["top", "middle", "bottom"]:
        def _build(prs, _s=anchor_slug, _xml=_ANCHOR_XML[anchor_slug]):
            sld = prs.slides.add_slide(prs.slide_layouts[6])
            shp = sld.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _emu(2), _emu(1), _emu(6), _emu(4))
            tf = shp.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"Anchor: {_s}"
            run.font.size = Pt(24)
            run.font.name = "Calibri"
            run.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            bodyPr = tf._txBody.find(qn("a:bodyPr"))
            bodyPr.set("anchor", _xml)
        _add(f"anchor-{anchor_slug}", _build)

    # --- Line spacing ---
    def _build_spacing(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = sld.shapes.add_textbox(_emu(1), _emu(0.5), _emu(10), _emu(6))
        tf = txbox.text_frame
        tf.word_wrap = True
        for i in range(5):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"Line {i+1}: Lorem ipsum dolor sit amet, consectetur adipiscing elit."
            p.font.size = Pt(18)
            p.font.name = "Calibri"
            p.space_after = Pt(12)
            p.space_before = Pt(6)
    _add("line-spacing", _build_spacing)

    # --- Placeholder idx-only inheritance (slide ph idx -> layout/master ph type) ---
    def _build_placeholder_idx_inheritance(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content

        sld.shapes.title.text = "Placeholder Inheritance Title"

        body = sld.placeholders[1]
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = "Body placeholder inherits bullet and master text size"

        p2 = tf.add_paragraph()
        p2.text = "Second line also inherits body placeholder style"
        p2.level = 0
    _add(
        "placeholder-idx-inheritance",
        _build_placeholder_idx_inheritance,
        _patch_placeholder_idx_inheritance_case,
    )

    return cases


# ---------------------------------------------------------------------------
# P1: Shape adjustment variants
# ---------------------------------------------------------------------------

def _build_shape_adj_cases() -> list[CaseDef]:
    cases: list[CaseDef] = []
    seq = 0

    def _add(slug: str, build_fn):
        nonlocal seq
        seq += 1
        cases.append({
            "name": f"oracle-pypptx-shape-adj-{seq:04d}-{slug}",
            "build_fn": build_fn,
        })

    # Shapes with meaningful adjustments: (MSO_SHAPE, slug, [(adj_index, value), ...])
    adj_configs = [
        # roundRect: corner radius
        (MSO_SHAPE.ROUNDED_RECTANGLE, "round-rect-small-radius", [(0, 0.05)]),
        (MSO_SHAPE.ROUNDED_RECTANGLE, "round-rect-large-radius", [(0, 0.45)]),
        # chevron: point depth
        (MSO_SHAPE.CHEVRON, "chevron-shallow", [(0, 0.15)]),
        (MSO_SHAPE.CHEVRON, "chevron-deep", [(0, 0.45)]),
        # right arrow: head width and depth
        (MSO_SHAPE.RIGHT_ARROW, "arrow-thin", [(0, 0.2), (1, 0.3)]),
        (MSO_SHAPE.RIGHT_ARROW, "arrow-wide-head", [(0, 0.1), (1, 0.6)]),
        # star 5-point: inner radius
        (MSO_SHAPE.STAR_5_POINT, "star5-thin", [(0, 0.15)]),
        (MSO_SHAPE.STAR_5_POINT, "star5-fat", [(0, 0.45)]),
        # donut: ring thickness
        (MSO_SHAPE.DONUT, "donut-thin-ring", [(0, 0.1)]),
        (MSO_SHAPE.DONUT, "donut-thick-ring", [(0, 0.45)]),
        # cross: arm thickness
        (MSO_SHAPE.CROSS, "cross-thin", [(0, 0.15)]),
        (MSO_SHAPE.CROSS, "cross-thick", [(0, 0.45)]),
        # trapezoid
        (MSO_SHAPE.TRAPEZOID, "trapezoid-narrow-top", [(0, 0.15)]),
        (MSO_SHAPE.TRAPEZOID, "trapezoid-wide-top", [(0, 0.45)]),
        # block arc
        (MSO_SHAPE.BLOCK_ARC, "block-arc-narrow", [(0, 0.1)]),
        (MSO_SHAPE.BLOCK_ARC, "block-arc-wide", [(0, 0.4)]),
        # folded corner
        (MSO_SHAPE.FOLDED_CORNER, "folded-corner-small", [(0, 0.1)]),
        (MSO_SHAPE.FOLDED_CORNER, "folded-corner-large", [(0, 0.4)]),
        # bevel
        (MSO_SHAPE.BEVEL, "bevel-thin", [(0, 0.05)]),
        (MSO_SHAPE.BEVEL, "bevel-thick", [(0, 0.35)]),
        # isosceles triangle: peak offset
        (MSO_SHAPE.ISOSCELES_TRIANGLE, "triangle-left-peak", [(0, 0.1)]),
        (MSO_SHAPE.ISOSCELES_TRIANGLE, "triangle-right-peak", [(0, 0.9)]),
        # pentagon
        (MSO_SHAPE.PENTAGON, "pentagon-shallow", [(0, 0.15)]),
        (MSO_SHAPE.PENTAGON, "pentagon-deep", [(0, 0.45)]),
        # can (cylinder): top ellipse height
        (MSO_SHAPE.CAN, "can-flat-top", [(0, 0.1)]),
        (MSO_SHAPE.CAN, "can-tall-top", [(0, 0.4)]),
        # heart
        (MSO_SHAPE.HEART, "heart-default", []),
        # moon
        (MSO_SHAPE.MOON, "moon-thin-crescent", [(0, 0.15)]),
        (MSO_SHAPE.MOON, "moon-wide-crescent", [(0, 0.7)]),
        # left brace
        (MSO_SHAPE.LEFT_BRACE, "left-brace-sharp", [(0, 0.05)]),
        (MSO_SHAPE.LEFT_BRACE, "left-brace-round", [(0, 0.3)]),
    ]

    for shape_type, slug, adjs in adj_configs:
        def _build(prs, _st=shape_type, _adjs=adjs):
            sld = prs.slides.add_slide(prs.slide_layouts[6])
            shp = sld.shapes.add_shape(_st, _emu(2), _emu(1), _emu(5), _emu(4))
            for idx, val in _adjs:
                try:
                    shp.adjustments[idx] = val
                except (IndexError, ValueError) as exc:
                    print(f"    WARN: adj[{idx}]={val} failed on {_st}: {exc}", flush=True)
        _add(slug, _build)

    return cases


# ---------------------------------------------------------------------------
# P2: Composite (multi-component) cases
# ---------------------------------------------------------------------------

def _build_composite_cases() -> list[CaseDef]:
    cases: list[CaseDef] = []
    seq = 0

    def _add(slug: str, build_fn):
        nonlocal seq
        seq += 1
        cases.append({
            "name": f"oracle-pypptx-composite-{seq:04d}-{slug}",
            "build_fn": build_fn,
        })

    # --- Two shapes side by side ---
    def _build_two_shapes(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        sld.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _emu(1), _emu(1.5), _emu(4), _emu(3))
        sld.shapes.add_shape(MSO_SHAPE.OVAL, _emu(6.5), _emu(1.5), _emu(4), _emu(3))
    _add("two-shapes-side-by-side", _build_two_shapes)

    # --- Shape with text inside ---
    def _build_shape_with_text(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        shp = sld.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _emu(2), _emu(1), _emu(6), _emu(4))
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Text inside a rounded rectangle"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.name = "Calibri"
        p.font.bold = True
    _add("shape-with-centered-text", _build_shape_with_text)

    # --- Shape + Textbox overlay ---
    def _build_shape_textbox_overlay(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        sld.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(1.5), _emu(1), _emu(7), _emu(5))
        txbox = sld.shapes.add_textbox(_emu(2), _emu(2), _emu(6), _emu(3))
        tf = txbox.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "Overlaid Title"
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(32)
        p1.font.bold = True
        p2 = tf.add_paragraph()
        p2.text = "Body text overlaid on a rectangle shape"
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(18)
    _add("shape-textbox-overlay", _build_shape_textbox_overlay)

    # --- Multiple shapes (grid) ---
    def _build_shape_grid(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = [MSO_SHAPE.RECTANGLE, MSO_SHAPE.OVAL, MSO_SHAPE.DIAMOND,
                  MSO_SHAPE.HEXAGON, MSO_SHAPE.STAR_5_POINT, MSO_SHAPE.HEART,
                  MSO_SHAPE.CROSS, MSO_SHAPE.RIGHT_ARROW, MSO_SHAPE.DONUT]
        for i, st in enumerate(shapes):
            row, col = divmod(i, 3)
            left = _emu(1 + col * 3.5)
            top = _emu(0.5 + row * 2.2)
            sld.shapes.add_shape(st, left, top, _emu(2.5), _emu(1.8))
    _add("shape-grid-3x3", _build_shape_grid)

    # --- Table + textbox ---
    def _build_table_textbox(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        # Title textbox
        txbox = sld.shapes.add_textbox(_emu(1), _emu(0.3), _emu(8), _emu(0.8))
        p = txbox.text_frame.paragraphs[0]
        p.text = "Quarterly Results"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.name = "Calibri"
        # Table
        tbl_shape = sld.shapes.add_table(4, 4, _emu(1), _emu(1.5), _emu(8), _emu(4))
        tbl = tbl_shape.table
        headers = ["Quarter", "Revenue", "Cost", "Profit"]
        data = [
            ["Q1", "$120K", "$85K", "$35K"],
            ["Q2", "$145K", "$92K", "$53K"],
            ["Q3", "$132K", "$88K", "$44K"],
        ]
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(14)
        for i, row in enumerate(data):
            for j, val in enumerate(row):
                tbl.cell(i + 1, j).text = val
    _add("table-with-title", _build_table_textbox)

    # --- Chart + title textbox ---
    def _build_chart_title(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        txbox = sld.shapes.add_textbox(_emu(1), _emu(0.3), _emu(8), _emu(0.8))
        p = txbox.text_frame.paragraphs[0]
        p.text = "Sales Overview"
        p.font.size = Pt(28)
        p.font.bold = True
        chart_data = CategoryChartData()
        chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
        chart_data.add_series("Product A", (45, 52, 48, 61))
        chart_data.add_series("Product B", (32, 38, 41, 35))
        sld.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            _emu(1), _emu(1.5), _emu(8), _emu(5),
            chart_data,
        )
    _add("chart-with-title", _build_chart_title)

    # --- Multiple textboxes with different styles ---
    def _build_multi_text(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        configs = [
            (_emu(0.5), _emu(0.5), _emu(4), _emu(2), "Arial", 20, True, PP_ALIGN.LEFT, "Left-aligned Arial Bold"),
            (_emu(5.5), _emu(0.5), _emu(5), _emu(2), "Georgia", 20, False, PP_ALIGN.RIGHT, "Right-aligned Georgia Italic"),
            (_emu(0.5), _emu(3), _emu(10), _emu(2), "Calibri", 16, False, PP_ALIGN.CENTER, "Centered Calibri — Lorem ipsum dolor sit amet, consectetur adipiscing elit."),
            (_emu(0.5), _emu(5.5), _emu(10), _emu(1.5), "Courier New", 14, False, PP_ALIGN.LEFT, "Monospace: code_sample = function(x) { return x * 2; }"),
        ]
        for left, top, w, h, font, size, bold, align, text in configs:
            txbox = sld.shapes.add_textbox(left, top, w, h)
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = align
            p.font.name = font
            p.font.size = Pt(size)
            p.font.bold = bold
            if font == "Georgia":
                p.font.italic = True
    _add("multi-textbox-styles", _build_multi_text)

    # --- Two charts side by side ---
    def _build_two_charts(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd1 = CategoryChartData()
        cd1.categories = ["Jan", "Feb", "Mar"]
        cd1.add_series("Sales", (120, 135, 148))
        sld.shapes.add_chart(XL_CHART_TYPE.LINE, _emu(0.5), _emu(0.5), _emu(5.5), _emu(6), cd1)

        cd2 = CategoryChartData()
        cd2.categories = ["A", "B", "C", "D"]
        cd2.add_series("Share", (35, 25, 22, 18))
        sld.shapes.add_chart(XL_CHART_TYPE.PIE, _emu(6.5), _emu(0.5), _emu(5.5), _emu(6), cd2)
    _add("two-charts-line-pie", _build_two_charts)

    # --- Shapes with different fills + text ---
    def _build_colored_shapes_text(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        configs = [
            (MSO_SHAPE.ROUNDED_RECTANGLE, _emu(0.5), _emu(1), RGBColor(0x00, 0x70, 0xC0), "Blue Box"),
            (MSO_SHAPE.OVAL, _emu(4.5), _emu(1), RGBColor(0xFF, 0x40, 0x40), "Red Oval"),
            (MSO_SHAPE.HEXAGON, _emu(8.5), _emu(1), RGBColor(0x00, 0xB0, 0x50), "Green Hex"),
        ]
        for st, left, top, color, text in configs:
            shp = sld.shapes.add_shape(st, left, top, _emu(3.5), _emu(4))
            shp.fill.solid()
            shp.fill.fore_color.rgb = color
            tf = shp.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _add("colored-shapes-with-text", _build_colored_shapes_text)

    # --- Table + chart + text ---
    def _build_dashboard(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        # Title
        txbox = sld.shapes.add_textbox(_emu(0.5), _emu(0.2), _emu(12), _emu(0.7))
        p = txbox.text_frame.paragraphs[0]
        p.text = "Dashboard: Monthly KPIs"
        p.font.size = Pt(28)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        # Small table (left)
        tbl_shape = sld.shapes.add_table(3, 2, _emu(0.5), _emu(1.2), _emu(4), _emu(2.5))
        tbl = tbl_shape.table
        for j, h in enumerate(["Metric", "Value"]):
            tbl.cell(0, j).text = h
        tbl.cell(1, 0).text = "Users"
        tbl.cell(1, 1).text = "12,450"
        tbl.cell(2, 0).text = "Revenue"
        tbl.cell(2, 1).text = "$89K"
        # Chart (right)
        cd = CategoryChartData()
        cd.categories = ["Mon", "Tue", "Wed", "Thu", "Fri"]
        cd.add_series("Visits", (850, 920, 780, 1050, 990))
        sld.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            _emu(5), _emu(1.2), _emu(7.5), _emu(5.5), cd,
        )
    _add("dashboard-table-chart", _build_dashboard)

    return cases


# ---------------------------------------------------------------------------
# P3: Chart data variants (2D types only — ECharts renderable)
# ---------------------------------------------------------------------------

def _build_chart_cases() -> list[CaseDef]:
    cases: list[CaseDef] = []
    seq = 0

    def _add(slug: str, build_fn):
        nonlocal seq
        seq += 1
        cases.append({
            "name": f"oracle-pypptx-chart-{seq:04d}-{slug}",
            "build_fn": build_fn,
        })

    # --- Column/Bar variants ---
    def _build_col_multi_series(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["Q1", "Q2", "Q3", "Q4"]
        cd.add_series("Product A", (45, 52, 48, 61))
        cd.add_series("Product B", (32, 38, 41, 35))
        cd.add_series("Product C", (28, 31, 36, 42))
        sld.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("column-clustered-3series", _build_col_multi_series)

    def _build_col_negative(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
        cd.add_series("Profit/Loss", (15, -8, 22, -12, 5, -3))
        sld.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("column-negative-values", _build_col_negative)

    def _build_col_stacked(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["2021", "2022", "2023", "2024"]
        cd.add_series("Hardware", (120, 135, 142, 158))
        cd.add_series("Software", (85, 102, 118, 131))
        cd.add_series("Services", (45, 52, 68, 79))
        sld.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("column-stacked-3series", _build_col_stacked)

    def _build_col_100_stacked(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["A", "B", "C"]
        cd.add_series("X", (30, 40, 25))
        cd.add_series("Y", (50, 35, 45))
        cd.add_series("Z", (20, 25, 30))
        sld.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED_100, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("column-100-stacked", _build_col_100_stacked)

    def _build_bar_clustered(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["Engineering", "Sales", "Marketing", "Support", "HR"]
        cd.add_series("Headcount", (45, 32, 18, 25, 8))
        sld.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("bar-clustered-single", _build_bar_clustered)

    def _build_bar_stacked(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["Dept A", "Dept B", "Dept C"]
        cd.add_series("FY23", (120, 95, 80))
        cd.add_series("FY24", (135, 110, 92))
        sld.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("bar-stacked-2series", _build_bar_stacked)

    # --- Line variants ---
    def _build_line_multi(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
        cd.add_series("Website", (1200, 1350, 1100, 1450, 1380, 1520))
        cd.add_series("Mobile", (800, 920, 850, 1050, 1100, 1180))
        cd.add_series("API", (300, 350, 380, 420, 460, 510))
        sld.shapes.add_chart(XL_CHART_TYPE.LINE, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("line-3series", _build_line_multi)

    def _build_line_markers(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["W1", "W2", "W3", "W4", "W5", "W6", "W7", "W8"]
        cd.add_series("Actual", (82, 85, 79, 91, 88, 94, 87, 96))
        cd.add_series("Target", (85, 85, 85, 85, 90, 90, 90, 90))
        sld.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("line-with-markers", _build_line_markers)

    def _build_line_stacked(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["Mon", "Tue", "Wed", "Thu", "Fri"]
        cd.add_series("Email", (120, 132, 101, 134, 90))
        cd.add_series("Chat", (220, 182, 191, 234, 290))
        sld.shapes.add_chart(XL_CHART_TYPE.LINE_STACKED, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("line-stacked", _build_line_stacked)

    # --- Pie / Doughnut ---
    def _build_pie_many(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["Chrome", "Safari", "Firefox", "Edge", "Other"]
        cd.add_series("Browser Share", (64, 19, 4, 5, 8))
        sld.shapes.add_chart(XL_CHART_TYPE.PIE, _emu(2), _emu(0.5), _emu(8), _emu(6), cd)
    _add("pie-5-categories", _build_pie_many)

    def _build_pie_exploded(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["A", "B", "C", "D"]
        cd.add_series("Sales", (40, 25, 20, 15))
        sld.shapes.add_chart(XL_CHART_TYPE.PIE_EXPLODED, _emu(2), _emu(0.5), _emu(8), _emu(6), cd)
    _add("pie-exploded", _build_pie_exploded)

    def _build_doughnut(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["Complete", "In Progress", "Not Started"]
        cd.add_series("Status", (65, 20, 15))
        sld.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, _emu(2), _emu(0.5), _emu(8), _emu(6), cd)
    _add("doughnut-3-categories", _build_doughnut)

    def _build_doughnut_exploded(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["A", "B", "C"]
        cd.add_series("Values", (50, 30, 20))
        sld.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT_EXPLODED, _emu(2), _emu(0.5), _emu(8), _emu(6), cd)
    _add("doughnut-exploded", _build_doughnut_exploded)

    # --- Area ---
    def _build_area(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["2020", "2021", "2022", "2023", "2024"]
        cd.add_series("Revenue", (80, 95, 110, 125, 148))
        cd.add_series("Cost", (60, 68, 75, 82, 91))
        sld.shapes.add_chart(XL_CHART_TYPE.AREA, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("area-2series", _build_area)

    def _build_area_stacked(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["Q1", "Q2", "Q3", "Q4"]
        cd.add_series("Product", (45, 52, 48, 55))
        cd.add_series("Service", (30, 35, 42, 38))
        cd.add_series("Support", (15, 18, 20, 22))
        sld.shapes.add_chart(XL_CHART_TYPE.AREA_STACKED, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("area-stacked-3series", _build_area_stacked)

    # --- Scatter ---
    def _build_scatter(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = XyChartData()
        s1 = cd.add_series("Cluster A")
        for x, y in [(1.2, 3.1), (2.4, 4.2), (3.1, 2.8), (1.8, 3.6), (2.9, 4.8)]:
            s1.add_data_point(x, y)
        s2 = cd.add_series("Cluster B")
        for x, y in [(5.1, 1.2), (6.3, 2.1), (5.8, 1.8), (7.1, 2.5), (6.0, 0.9)]:
            s2.add_data_point(x, y)
        sld.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("scatter-2-clusters", _build_scatter)

    def _build_scatter_smooth(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = XyChartData()
        s = cd.add_series("Curve")
        for i in range(20):
            x = i * 0.5
            y = math.sin(x) * 3 + 5
            s.add_data_point(x, y)
        sld.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER_SMOOTH, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("scatter-smooth-sine", _build_scatter_smooth)

    # --- Radar ---
    def _build_radar(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["Speed", "Power", "Agility", "Defense", "Stamina"]
        cd.add_series("Player A", (85, 70, 90, 65, 75))
        cd.add_series("Player B", (70, 85, 65, 80, 90))
        sld.shapes.add_chart(XL_CHART_TYPE.RADAR, _emu(2), _emu(0.5), _emu(8), _emu(6), cd)
    _add("radar-2series", _build_radar)

    def _build_radar_filled(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = ["Math", "Science", "English", "History", "Art"]
        cd.add_series("Student", (92, 85, 78, 88, 95))
        sld.shapes.add_chart(XL_CHART_TYPE.RADAR_FILLED, _emu(2), _emu(0.5), _emu(8), _emu(6), cd)
    _add("radar-filled-single", _build_radar_filled)

    # --- Bubble ---
    def _build_bubble(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = BubbleChartData()
        s = cd.add_series("Markets")
        s.add_data_point(1.5, 2.5, 10)
        s.add_data_point(3.0, 4.0, 25)
        s.add_data_point(5.0, 1.5, 15)
        s.add_data_point(2.5, 3.5, 30)
        sld.shapes.add_chart(XL_CHART_TYPE.BUBBLE, _emu(1), _emu(0.5), _emu(10), _emu(6), cd)
    _add("bubble-4-points", _build_bubble)

    # --- Large dataset ---
    def _build_line_large(prs):
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        cd = CategoryChartData()
        cd.categories = [str(i) for i in range(1, 25)]
        random.seed(42)
        base = 100
        vals = []
        for _ in range(24):
            base += random.randint(-10, 15)
            vals.append(base)
        cd.add_series("Monthly Trend", vals)
        sld.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, _emu(0.5), _emu(0.5), _emu(11), _emu(6), cd)
    _add("line-24-month-trend", _build_line_large)

    return cases


# ---------------------------------------------------------------------------
# Orchestrator
# ---------------------------------------------------------------------------

def _build_all_case_defs() -> list[CaseDef]:
    all_cases: list[CaseDef] = []
    all_cases.extend(_build_text_cases())
    all_cases.extend(_build_shape_adj_cases())
    all_cases.extend(_build_composite_cases())
    all_cases.extend(_build_chart_cases())
    return all_cases


def _generate_pptx(case_def: CaseDef, output_path: str | Path) -> None:
    """Generate a single PPTX file using python-pptx."""
    output_path = Path(output_path)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    case_def["build_fn"](prs)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    if postprocess_fn := case_def.get("postprocess_fn"):
        postprocess_fn(output_path)


def _write_case_json(case_def: CaseDef, cases_dir: Path) -> Path:
    """Write a minimal case JSON for eval script discovery."""
    name = case_def["name"]
    payload = {
        "name": name,
        "generator": "python-pptx",
        "slides": [{"nodes": [{"kind": "pypptx-generated"}]}],
    }
    out = cases_dir / f"{name}.json"
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(payload, indent=2) + "\n", encoding="utf-8")
    return out


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Generate ground truth cases using python-pptx + PowerPoint PDF export.",
    )
    parser.add_argument("--cases-dir", type=Path, default=CASES_DIR)
    parser.add_argument("--testdata-dir", type=Path, default=TESTDATA_DIR)
    parser.add_argument("--report-path", type=Path, default=REPORT_PATH)
    parser.add_argument("--pptx-only", action="store_true",
                        help="Only generate PPTX files (skip PDF export). Works on any platform.")
    parser.add_argument("--export-png", action="store_true", default=True,
                        help="Export each slide as PNG (default: enabled, Windows only).")
    parser.add_argument("--no-export-png", action="store_false", dest="export_png",
                        help="Skip PNG export.")
    parser.add_argument("--png-width", type=int, default=0, metavar="PX",
                        help="PNG export width in pixels (0 = PowerPoint default).")
    parser.add_argument("--png-height", type=int, default=0, metavar="PX",
                        help="PNG export height in pixels (0 = PowerPoint default).")
    parser.add_argument("--no-reuse", action="store_true",
                        help="Force regeneration even if files exist.")
    args = parser.parse_args()

    cases_dir = args.cases_dir.resolve()
    testdata_dir = args.testdata_dir.resolve()
    report_path = args.report_path.resolve()

    all_cases = _build_all_case_defs()
    print(f"Total case definitions: {len(all_cases)}")

    generated: list[str] = []
    failures: list[dict] = []
    skipped = 0

    do_png = args.export_png and not args.pptx_only

    # Import ground truth export only if needed (per-case COM session for stability)
    export_fn = None
    if not args.pptx_only:
        if sys.platform != "win32":
            print("ERROR: PDF/PNG export requires Windows + PowerPoint. Use --pptx-only on other platforms.",
                  file=sys.stderr)
            return 1
        from oracle.powerpoint_oracle import export_pptx_ground_truth_win
        export_fn = export_pptx_ground_truth_win

    for i, case_def in enumerate(all_cases, 1):
        name = case_def["name"]
        case_d = testdata_dir / "cases" / name
        pptx_path = case_d / "source.pptx"
        pdf_path = case_d / "ground-truth.pdf"
        slides_d = case_d / "slides"

        # Reuse check
        if not args.no_reuse:
            if args.pptx_only and pptx_path.exists():
                skipped += 1
                continue
            if not args.pptx_only and pptx_path.exists() and pdf_path.exists():
                # If PNG export requested but slide1.png missing, regenerate
                if do_png and not (slides_d / "slide1.png").exists():
                    pass  # fall through
                else:
                    skipped += 1
                    continue

        print(f"  [{i}/{len(all_cases)}] {name} ...", end=" ", flush=True)

        try:
            # Generate PPTX via python-pptx
            _generate_pptx(case_def, pptx_path)

            # Export PDF + PNG via independent PowerPoint COM session (one per case)
            if export_fn is not None:
                export_fn(
                    pptx_path, pdf_path,
                    slides_png_dir=slides_d if do_png else None,
                    png_width=args.png_width,
                    png_height=args.png_height,
                )

            # Write case JSON for eval discovery
            _write_case_json(case_def, cases_dir)

            generated.append(name)
            print("OK")
        except Exception as exc:
            failures.append({"case": name, "error": str(exc)})
            print(f"FAIL: {exc}")

    print(f"\nResults: {len(generated)} generated, {skipped} reused, {len(failures)} failed")

    report = {
        "generator": "python-pptx",
        "cases_dir": str(cases_dir),
        "testdata_dir": str(testdata_dir),
        "pptx_only": args.pptx_only,
        "export_png": do_png,
        "png_width": args.png_width,
        "png_height": args.png_height,
        "total_definitions": len(all_cases),
        "generated_count": len(generated),
        "skipped_reused": skipped,
        "failed_count": len(failures),
        "generated_cases": generated,
        "failed_cases": failures,
    }
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_text(json.dumps(report, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    print(f"Report: {report_path}")

    return 1 if failures else 0


if __name__ == "__main__":
    sys.exit(main())
